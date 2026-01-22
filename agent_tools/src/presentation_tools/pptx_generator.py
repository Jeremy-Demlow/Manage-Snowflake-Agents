"""
PowerPoint Generator for SI Conversations
Creates professional presentations from agent Q&A data with tables and charts.
"""

import re
import json
import tempfile
from datetime import datetime
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData


BRAND_COLORS = {
    "primary": RGBColor(0x1A, 0x73, 0xE8),
    "secondary": RGBColor(0x0D, 0x47, 0xA1),
    "accent": RGBColor(0x42, 0x85, 0xF4),
    "dark": RGBColor(0x33, 0x33, 0x33),
    "light": RGBColor(0xF5, 0xF5, 0xF5),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}


def create_conversation_pptx(
    question: str,
    response: str,
    title: Optional[str] = None,
    subtitle: Optional[str] = None,
) -> str:
    """
    Create a PowerPoint presentation from an SI conversation.

    Args:
        question: User's original question
        response: Agent's markdown response
        title: Optional custom title (defaults to truncated question)
        subtitle: Optional subtitle

    Returns:
        Path to generated .pptx file
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    display_title = title or _truncate(question, 60)
    display_subtitle = subtitle or f"Generated {datetime.now().strftime('%B %d, %Y')}"

    _add_title_slide(prs, display_title, display_subtitle)
    _add_executive_summary_slide(prs, question, response)
    _add_qa_slides(prs, question, response)

    tables = _extract_tables(response)
    for i, table_data in enumerate(tables[:3]):
        _add_table_slide(prs, table_data, f"Data Table {i+1}")

    chart_data = _extract_chart_data(response)
    if chart_data:
        _add_chart_slide(prs, chart_data)

    _add_closing_slide(prs)

    with tempfile.NamedTemporaryFile(
        mode="wb", prefix="si_report_", suffix=".pptx", delete=False
    ) as f:
        prs.save(f.name)
        return f.name


def _add_title_slide(prs: Presentation, title: str, subtitle: str):
    """Add title slide with gradient-style header."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    header_shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(3)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = BRAND_COLORS["primary"]
    header_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(1), Inches(11.833), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BRAND_COLORS["white"]
    p.alignment = PP_ALIGN.LEFT

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(3.5), Inches(11.833), Inches(1)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = BRAND_COLORS["dark"]
    p.alignment = PP_ALIGN.LEFT

    footer_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(6.5), Inches(11.833), Inches(0.5)
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Powered by Snowflake Intelligence"
    p.font.size = Pt(14)
    p.font.color.rgb = BRAND_COLORS["accent"]


def _add_executive_summary_slide(prs: Presentation, question: str, response: str):
    """Add executive summary slide with key insights."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_slide_header(slide, "Executive Summary")

    summary = _generate_summary(response)

    content_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Question Asked:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BRAND_COLORS["dark"]

    p = tf.add_paragraph()
    p.text = _truncate(question, 200)
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = BRAND_COLORS["secondary"]
    p.space_after = Pt(20)

    p = tf.add_paragraph()
    p.text = "Key Insights:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BRAND_COLORS["dark"]
    p.space_before = Pt(10)

    for point in summary[:5]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.font.color.rgb = BRAND_COLORS["dark"]
        p.level = 0


def _add_qa_slides(prs: Presentation, question: str, response: str):
    """Add Q&A content slides, splitting long responses."""
    sections = _split_response(response)

    for i, section in enumerate(sections[:4]):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        section_title = section.get("title", f"Response {i+1}")
        _add_slide_header(slide, section_title)

        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        content = section.get("content", "")
        lines = content.split("\n")

        for j, line in enumerate(lines[:15]):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            clean_line = line.strip()
            if clean_line.startswith(("- ", "* ", "• ")):
                p.text = f"• {clean_line[2:]}"
                p.level = 0
            elif clean_line.startswith(("1.", "2.", "3.", "4.", "5.")):
                p.text = clean_line
                p.font.bold = True
            else:
                p.text = clean_line

            p.font.size = Pt(14)
            p.font.color.rgb = BRAND_COLORS["dark"]


def _add_table_slide(prs: Presentation, table_data: dict, title: str):
    """Add a slide with a data table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_slide_header(slide, title)

    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])

    if not headers or not rows:
        return

    num_cols = len(headers)
    num_rows = min(len(rows) + 1, 10)

    table_width = min(Inches(12), Inches(2 * num_cols))
    left = (prs.slide_width - table_width) / 2

    table = slide.shapes.add_table(
        num_rows, num_cols, left, Inches(1.5), table_width, Inches(0.5 * num_rows)
    ).table

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_COLORS["primary"]
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = BRAND_COLORS["white"]
        p.alignment = PP_ALIGN.CENTER

    for row_idx, row in enumerate(rows[: num_rows - 1]):
        for col_idx, value in enumerate(row[:num_cols]):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = BRAND_COLORS["dark"]
            p.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BRAND_COLORS["light"]


def _add_chart_slide(prs: Presentation, chart_data: dict):
    """Add a slide with a bar chart."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    chart_title = chart_data.get("title", "Data Visualization")
    _add_slide_header(slide, chart_title)

    cd = CategoryChartData()
    cd.categories = chart_data.get("categories", [])

    for series in chart_data.get("series", []):
        cd.add_series(series.get("name", "Values"), series.get("values", []))

    x, y, cx, cy = Inches(1), Inches(1.5), Inches(11.333), Inches(5.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, cd
    ).chart

    chart.has_legend = len(chart_data.get("series", [])) > 1


def _add_closing_slide(prs: Presentation):
    """Add closing/thank you slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg_shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BRAND_COLORS["primary"]
    bg_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(2.5), Inches(11.833), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BRAND_COLORS["white"]
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.5), Inches(11.833), Inches(1)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Generated by Snowflake Intelligence"
    p.font.size = Pt(24)
    p.font.color.rgb = BRAND_COLORS["white"]
    p.alignment = PP_ALIGN.CENTER


def _add_slide_header(slide, title: str):
    """Add consistent header bar to content slides."""
    header = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = BRAND_COLORS["primary"]
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRAND_COLORS["white"]


def _truncate(text: str, max_len: int) -> str:
    """Truncate text with ellipsis."""
    return text[: max_len - 3] + "..." if len(text) > max_len else text


def _generate_summary(response: str) -> list:
    """Extract key points from response for executive summary."""
    points = []

    lines = response.split("\n")
    for line in lines:
        line = line.strip()
        if line.startswith(("- ", "* ", "• ", "1.", "2.", "3.")):
            clean = re.sub(r"^[-*•\d.]+\s*", "", line)
            if 10 < len(clean) < 200:
                points.append(clean)

    if len(points) < 3:
        sentences = re.split(r"[.!?]+", response)
        for s in sentences:
            s = s.strip()
            if 20 < len(s) < 200 and s not in points:
                points.append(s)
                if len(points) >= 5:
                    break

    return points[:5]


def _split_response(response: str) -> list:
    """Split response into logical sections."""
    sections = []

    header_pattern = r"^(#{1,3})\s+(.+)$"
    current_section = {"title": "Overview", "content": ""}

    for line in response.split("\n"):
        match = re.match(header_pattern, line)
        if match:
            if current_section["content"].strip():
                sections.append(current_section)
            current_section = {"title": match.group(2), "content": ""}
        else:
            current_section["content"] += line + "\n"

    if current_section["content"].strip():
        sections.append(current_section)

    if not sections:
        sections.append({"title": "Response", "content": response})

    return sections


def _extract_tables(response: str) -> list:
    """Extract markdown tables from response."""
    tables = []

    table_pattern = r"\|(.+)\|\n\|[-:\s|]+\|\n((?:\|.+\|\n?)+)"
    matches = re.findall(table_pattern, response)

    for header_row, body in matches:
        headers = [h.strip() for h in header_row.split("|") if h.strip()]
        rows = []
        for row_line in body.strip().split("\n"):
            cells = [c.strip() for c in row_line.split("|") if c.strip()]
            if cells:
                rows.append(cells)

        if headers and rows:
            tables.append({"headers": headers, "rows": rows})

    return tables


def _extract_chart_data(response: str) -> Optional[dict]:
    """Attempt to extract chartable data from response."""
    tables = _extract_tables(response)

    for table in tables:
        headers = table.get("headers", [])
        rows = table.get("rows", [])

        if len(headers) >= 2 and len(rows) >= 2:
            categories = []
            values = []

            for row in rows[:10]:
                if len(row) >= 2:
                    categories.append(str(row[0]))
                    try:
                        val = float(re.sub(r"[,$%]", "", str(row[1])))
                        values.append(val)
                    except ValueError:
                        continue

            if len(categories) >= 2 and len(values) == len(categories):
                return {
                    "title": headers[1] if len(headers) > 1 else "Data",
                    "categories": categories,
                    "series": [{"name": headers[1], "values": values}],
                }

    return None
